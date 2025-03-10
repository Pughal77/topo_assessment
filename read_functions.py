'''
Read functions for diff types. For visualisation of data we are to use matplot lib
'''
import pandas as pd
import tabula
from pptx import Presentation
import json
'''
    This class represents the unified data structure I am to return
    Upon looking at the datasets, dataset 1 is about 2 sport and fitness companies
    dataset 3 is about quarterly revenue n stuff
    datset 4 represents the summary of performance
    for data visualisation, I think we could go about verifying the summaries for fitpro
    The unified data structure although it has 
'''
class Unified_data_structure:
    def __init__(self):
        self.data = {
            "dataset_1": None,
            "dataset_2": None,
            "dataset_3": None,
            "dataset_4": None
        }
        self.read_json("datasets/dataset1.json")
        self.read_csv("datasets/dataset2.csv")
        self.read_pdf("datasets/dataset3.pdf")
        self.read_pptx("datasets/dataset4.pptx")

    def read_json(self, filename) -> None:
        # data might be nested
        json_data = pd.read_json(filename)
        self.data['dataset_1'] = json_data
        # after exploring the json_data object in pyshell
        # i think it is more meaningful to structure this data as such
        # for now, we replace the list of employees to a pandas dataframe as well
        headers = [key for key, list in json_data["companies"][0]["employees"][0].items()]
        # replace it as pandas dataframe
        for i in range(len(json_data['companies'])):
            json_data["companies"][i]["employees"] = pd.DataFrame(json_data["companies"][i]["employees"], columns=headers)


    def read_csv(self, filename) -> None:
        csv_data = pd.read_csv(filename)
        self.data['dataset_2'] = csv_data
    
    def read_pdf(self, filename) -> None:
        # use tabula
        # read and injest data from pdf
        dfs = tabula.read_pdf(filename, pages='all')
        # from checking what dfs returns, I see that it is a list with a single element
        # where dfs[0] is a pandas dataframe
        self.data['dataset_3'] = dfs[0]
        return None
    def read_pptx(self, filename) -> None:
        # convert to pdf then read pdf
        # will convert all the pptx files to pdf in specified dir
        # in this dataset we have a table in slide 2
        pptx = Presentation(filename)
        # after exploring the pptx object in python shell i found where table in slide 2 is stored
        table = pptx.slides[1].shapes[1].table
        # convert this table to a pandas dataframe
        data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            data.append(row_data)
        
        # Convert to Pandas DataFrame
        headers = data[0]  # first row does contain headers
        rows = data[1:]  # Remaining rows are data
        df = pd.DataFrame(rows, columns=headers)
        
        # as for the information in slide 1 and 3, I think its best if we manually keyed in the information
        self.data['dataset_4'] = df

    def get_data(self):
        # return entire self.data as JSON
        json_data = {key: df.to_dict(orient='records') for key, df in self.data.items()}

        # Save the JSON data to a file
        with open('datasets/consolidated_datasets.json', 'w') as json_file:
            json.dump(json_data, json_file, indent=4)

    def get_data_xlsx(self):
        # here we are converting the self.data to an xlsx file 
        # in which each dataset would have its own excel sheet
        # for dataset 1, I think it is more meaningful to separate them into 2 sheets, one containing 
        with pd.ExcelWriter('datasets/consolidated_dataset.xlsx', engine='xlsxwriter') as writer:
        # Iterate through the dictionary and write each DataFrame to a separate sheet
            for sheet_name, df in self.data.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)
    
    def visualise(self):
        return None

x = Unified_data_structure()