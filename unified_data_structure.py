import pandas as pd
import tabula
from pptx import Presentation
import json
import matplotlib.pyplot as plt

"""
This class represents the unified data structure that we will ingest the 4 datasets and provide useful visualisations of this dataset.
The unified dataset is of the following structure:
    self.data represents the unified dataset. To make its implementation simple, it would follow the following rules:
    1) Tables of data would be represented as a pandas dataframe (so that we can make meanigful changes to the data)
    2) self.data is a dictionary where every entry has to be a key-value pair
        key could is a string description of the value
        value could be either a pandas dataframe or another dictionary
    
Data ingestion (How I am to unify the 4 datasets):
    From the first dataset, I am to create 3 entries
        First is a list of companies as pandas dataframe (has no emplyee data here)
        Second is a dictionary in which
            the key is the index of the company (as an integer) in the first list and the value is the list of employees of that company
            as a pandas dataframe
    From second dataset, it looks to be a list of information about customers who come to Sports and Leisure
        It will just be a list of customers as pandas data frame
        key = "customers", value = df
    From third dataset, is a single table which would be represesnted as a pandas dataframe as well
        key = "quarterly_performance", value = df
    From fourth dataset, there will be 3 dictionaries each providing the relavant data summaries
        key = "quarterly_metrics", value = dict
        key = "key_highlights", value = dict
        key = "revenue_distribution", value = dict
"""
class Unified_data_structure:
    def __init__(self):
        self.data = {}
        self.read_json("datasets/dataset1.json")
        self.read_csv("datasets/dataset2.csv")
        self.read_pdf("datasets/dataset3.pdf")
        self.read_pptx("datasets/dataset4.pptx")

    def read_json(self, filename):
        json_data = pd.read_json(filename)
        # after exploring the json_data object in pyshell, I realise that employees are still a list upon conversion
        # for now, we replace the list of employees to a pandas dataframe of employees
        header_employees = [key for key, list in json_data["companies"][0]["employees"][0].items()]
        # create companies and companies dictionary
        employees = {}
        companies = {}
        # populate companies dictionary
        for i in range(len(json_data['companies'])):
            employees[i] = pd.DataFrame(json_data["companies"][i]["employees"], columns=header_employees)
            del json_data["companies"][i]["employees"]
        header_companies = [key for key, list in json_data["companies"][0].items()]
        for i in range(len(json_data['companies'])):
            companies[i] = pd.DataFrame(json_data["companies"][i], columns=header_companies)
        self.data["companies"] = companies
        self.data["employees"] = employees
        
    def read_csv(self, filename):
        csv_data = pd.read_csv(filename)
        self.data['customers'] = csv_data
    
    def read_pdf(self, filename):
        # use tabula
        # read and injest data from pdf
        dfs = tabula.read_pdf(filename, pages='all')
        # from checking what dfs returns, I see that it is a list with a single element
        # where dfs[0] is a pandas dataframe
        self.data['quarterly_performance'] = dfs[0]
        return None
    
    def read_pptx(self, filename):
        # convert to pdf then read pdf
        # will convert all the pptx files to pdf in specified dir
        # in this dataset we have a table in slide 2
        pptx = Presentation(filename)
        # after exploring the pptx object in python shell i found where table in slide 2 is stored
        table = pptx.slides[1].shapes[1].table
        # convert this table to a pandas dataframe
        data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            data.append(row_data)
        
        # Convert to Pandas DataFrame
        headers = data[0]  # first row does contain headers
        rows = data[1:]  # Remaining rows are data
        quarterly_metrics = pd.DataFrame(rows, columns=headers)
        
        # as for the information in slide 1 and 3, I think its best if we manually keyed in the information as the information is not in a tabular format and is short
        # if there are too many datapoints then, my approach for that is to parse through the text in slide 1 and 3
        key_highlights = {
            "Total Revenue" : "$10,400,000",
            "Total Memberships Sold" : "1520",
            "Top Location" : "Downtown"
        }

        revenue_distribution = {
            "Gym": "40%",
            "Pool": "25%,",
            "Tennis Court": "15%",
            "Personal Training": "20%"
        }
        self.data["revenue_distribution"] = revenue_distribution
        self.data["key_highlights" ] = key_highlights
        self.data["quarterly_metrics"] = quarterly_metrics

    def get_data(self):
        # return entire self.data as JSON
        json_data = {key: df.to_dict(orient='records') for key, df in self.data.items()}

        # Save the JSON data to a file
        with open('datasets/consolidated_datasets.json', 'w') as json_file:
            json.dump(json_data, json_file, indent=4)

    def get_data_xlsx(self):
        # here we are converting the self.data to an xlsx file 
        # in which each dataset would have its own excel sheet
        # for dataset 1, I think it is more meaningful to separate them into 2 sheets, one containing 
        with pd.ExcelWriter('datasets/consolidated_dataset.xlsx', engine='xlsxwriter') as writer:
        # Iterate through the dictionary and write each DataFrame to a separate sheet
            for sheet_name, df in self.data.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)
    
    def visualise_data(self):
        # here is a visalisation of top revenue earners by activity amongst basic membership types
        # basic filtering and ordering of data is allowed as most values in the self.data dictionary are pandas dataframe

        fig, axes = plt.subplots(1, 2, figsize=(12, 6))  # 1 row, 2 columns

        # first visualisation is a graph plot of revenue earned from customers against the amount of time spent in the facility
        value_counts = dict(self.data["customers"]["Location"].value_counts())
        axes[0].bar(value_counts.keys(), value_counts.values())
        axes[0].set_title('Overview of customers in each location')
        axes[0].set_xlabel('Location')
        axes[0].set_ylabel('Number of customers')
        

        # second visualisation is quarterly revenue against time
        x = range(len(self.data["quarterly_performance"]))
        y = self.data["quarterly_performance"]["Revenue (in $)"]
        axes[1].plot(x, y)
        axes[1].set_title('Quarterly Performance')
        axes[1].set_xlabel('Time')
        axes[1].set_ylabel('Revenue')

        plt.tight_layout()
        plt.show()

x = Unified_data_structure()

x.visualise_data()
